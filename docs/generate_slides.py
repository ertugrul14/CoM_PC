"""
Generate a clean PPTX deck for Steps 9-12 of the Melbourne pipeline.
Run from repo root: python docs/generate_slides.py
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

OUT_PATH = Path(__file__).parent / "pipeline_steps_9_12.pptx"

# ── Palette ──────────────────────────────────────────────────────────────────
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
BLACK      = RGBColor(0x14, 0x14, 0x14)
NAVY       = RGBColor(0x0F, 0x37, 0x73)   # headings / title
BLUE       = RGBColor(0x28, 0x5A, 0xA0)   # subheadings
BLUE_LIGHT = RGBColor(0x46, 0x6E, 0xC0)   # accents
GREY_DARK  = RGBColor(0x3C, 0x3C, 0x3C)   # body text
GREY_MID   = RGBColor(0x78, 0x78, 0x78)   # muted / captions
GREY_LIGHT = RGBColor(0xF0, 0xF2, 0xF5)   # slide background
GREY_BOX   = RGBColor(0xE4, 0xE8, 0xEE)   # box fill
CODE_BG    = RGBColor(0xF5, 0xF5, 0xF5)
CODE_FG    = RGBColor(0x28, 0x28, 0x28)
ACCENT_RED = RGBColor(0xC0, 0x39, 0x2B)   # warning/note highlights
ACCENT_GRN = RGBColor(0x1A, 0x7A, 0x4A)   # positive


# ── Helpers ──────────────────────────────────────────────────────────────────
W  = Inches(13.33)   # widescreen 16:9
H  = Inches(7.5)

def new_prs() -> Presentation:
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs


def blank_slide(prs: Presentation):
    blank_layout = prs.slide_layouts[6]   # totally blank
    return prs.slides.add_slide(blank_layout)


def bg(slide, color: RGBColor = WHITE):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def txbox(slide, text, x, y, w, h,
          size=18, bold=False, italic=False,
          color=BLACK, align=PP_ALIGN.LEFT,
          font="Calibri"):
    """Add a simple text box. Returns the run for further styling."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    p  = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size    = Pt(size)
    run.font.bold    = bold
    run.font.italic  = italic
    run.font.color.rgb = color
    run.font.name    = font
    return run


def rect(slide, x, y, w, h, fill: RGBColor, line: RGBColor = None, line_w=Pt(0)):
    from pptx.util import Pt as _Pt
    shape = slide.shapes.add_shape(
        1,   # MSO_SHAPE_TYPE.RECTANGLE
        x, y, w, h
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line:
        shape.line.color.rgb = line
        shape.line.width = line_w
    else:
        shape.line.fill.background()
    return shape


def top_bar(slide, step_label: str, step_color: RGBColor = NAVY):
    """Draw the top accent bar with step label."""
    rect(slide, 0, 0, W, Inches(0.12), step_color)
    txbox(slide, step_label,
          Inches(0.35), Inches(0.18),
          Inches(12), Inches(0.42),
          size=11, bold=True, color=step_color, font="Calibri")


def slide_title(slide, title: str, subtitle: str = "",
                title_color=NAVY, sub_color=GREY_MID):
    txbox(slide, title,
          Inches(0.55), Inches(0.22),
          Inches(12.2), Inches(0.72),
          size=28, bold=True, color=title_color, font="Calibri")
    if subtitle:
        txbox(slide, subtitle,
              Inches(0.55), Inches(0.95),
              Inches(12.2), Inches(0.45),
              size=14, italic=True, color=sub_color, font="Calibri")


def section_head(slide, text, x, y, w=Inches(5.8), color=BLUE):
    txbox(slide, text, x, y, w, Inches(0.38),
          size=13, bold=True, color=color, font="Calibri")
    # underline via thin rect
    rect(slide, x, y + Inches(0.34), w, Inches(0.025), color)


def bullet_lines(slide, items: list[str], x, y, w, h,
                 size=11, color=GREY_DARK, indent="  ", spacing=1.2):
    """Add a multi-line bullet textbox."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.space_after  = Pt(spacing)
        p.space_before = Pt(0)
        run = p.add_run()
        run.text = item
        run.font.size  = Pt(size)
        run.font.color.rgb = color
        run.font.name  = "Calibri"


def code_box(slide, lines: list[str], x, y, w, h):
    """Monospace code block with light grey background."""
    rect(slide, x, y, w, h, CODE_BG, GREY_MID, Pt(0.5))
    text = "\n".join(lines)
    tb = slide.shapes.add_textbox(
        x + Inches(0.12), y + Inches(0.1),
        w - Inches(0.24), h - Inches(0.2)
    )
    tf = tb.text_frame
    tf.word_wrap = False
    first = True
    for line in lines:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        run = p.add_run()
        run.text = line
        run.font.size  = Pt(9.5)
        run.font.color.rgb = CODE_FG
        run.font.name  = "Courier New"


def info_box(slide, label, value, x, y, w=Inches(2.8), h=Inches(0.75),
             label_color=BLUE, val_color=BLACK):
    rect(slide, x, y, w, h, GREY_BOX)
    txbox(slide, label, x + Inches(0.12), y + Inches(0.06),
          w - Inches(0.2), Inches(0.28),
          size=9, bold=True, color=label_color)
    txbox(slide, value, x + Inches(0.12), y + Inches(0.34),
          w - Inches(0.2), Inches(0.34),
          size=11, bold=False, color=val_color)


def footer_note(slide, text):
    txbox(slide, text,
          Inches(0.35), Inches(7.1),
          Inches(12.6), Inches(0.3),
          size=8.5, italic=True, color=GREY_MID)


def divider(slide, y, color=GREY_BOX):
    rect(slide, Inches(0.35), y, Inches(12.6), Inches(0.02), color)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDES
# ═══════════════════════════════════════════════════════════════════════════════

def s_cover(prs):
    """Cover slide."""
    s = blank_slide(prs)
    bg(s, WHITE)
    # Left accent strip
    rect(s, 0, 0, Inches(0.45), H, NAVY)
    # Title area background
    rect(s, Inches(0.45), Inches(1.6), Inches(12.88), Inches(2.8), GREY_LIGHT)
    txbox(s, "Melbourne CBD Curbside Reallocation Pipeline",
          Inches(0.8), Inches(1.78), Inches(11.8), Inches(1.1),
          size=34, bold=True, color=NAVY, font="Calibri")
    txbox(s, "Steps 9 - 12: Training, Interpretation, Scenario Simulation & Export",
          Inches(0.8), Inches(2.88), Inches(11.8), Inches(0.6),
          size=17, italic=True, color=BLUE, font="Calibri")
    txbox(s, "Master's Thesis  |  Spatio-Temporal GNN  |  Melbourne CBD  |  2026",
          Inches(0.8), Inches(4.7), Inches(11.8), Inches(0.4),
          size=12, color=GREY_MID, font="Calibri")
    # Step badges
    for i, (num, label, col) in enumerate([
        ("09", "TRAIN",    NAVY),
        ("10", "INTERPRET",BLUE),
        ("11", "SCENARIO", BLUE_LIGHT),
        ("12", "EXPORT",   GREY_MID),
    ]):
        bx = Inches(0.8 + i * 3.1)
        rect(s, bx, Inches(5.35), Inches(2.6), Inches(0.88), col)
        txbox(s, f"STEP {num}", bx + Inches(0.12), Inches(5.42),
              Inches(2.4), Inches(0.35), size=10, bold=True, color=WHITE)
        txbox(s, label, bx + Inches(0.12), Inches(5.73),
              Inches(2.4), Inches(0.35), size=11, color=WHITE)


def s_overview(prs):
    """One-page overview: what each step does."""
    s = blank_slide(prs)
    bg(s, WHITE)
    top_bar(s, "OVERVIEW")
    slide_title(s, "Steps 9 - 12: What Each Step Does",
                "Post-training phase — model learns, interprets, simulates, and exports")

    rows = [
        ("09", "TRAIN",     NAVY,       "Trains the MultiGCN spatio-temporal GNN on the assembled data cube. Produces the trained model checkpoint used by all downstream steps."),
        ("10", "INTERPRET", BLUE,       "Quantifies which of the 23 input features the model actually relies on (permutation importance) and compares the contribution of spatial vs. semantic graph branches."),
        ("11", "SCENARIO",  BLUE_LIGHT, "Runs counterfactual rollouts: baseline vs. treated (e.g. pedestrianise a street). Reports pedestrian flow delta on the target street and network-wide spillover across all 1,397 streets."),
        ("12", "EXPORT",    GREY_DARK,  "Enriches GeoJSON with opportunity scores and writes frontend-ready JSON files for the Mapbox GL JS visualisation."),
    ]
    for i, (num, label, col, desc) in enumerate(rows):
        y = Inches(1.52 + i * 1.38)
        rect(s, Inches(0.35), y, Inches(1.55), Inches(1.1), col)
        txbox(s, f"STEP {num}", Inches(0.48), y + Inches(0.1),
              Inches(1.3), Inches(0.35), size=9, bold=True, color=WHITE)
        txbox(s, label, Inches(0.48), y + Inches(0.44),
              Inches(1.3), Inches(0.4), size=13, bold=True, color=WHITE)
        txbox(s, desc,
              Inches(2.1), y + Inches(0.12),
              Inches(10.8), Inches(0.88),
              size=12, color=GREY_DARK)


def s_step8(prs):
    """Step 08 — Data Cube Assembly."""
    s = blank_slide(prs)
    bg(s, WHITE)
    top_bar(s, "STEP 08  |  CUBE  |  data cube assembly + adjacency matrices", NAVY)
    slide_title(s, "Data Cube Assembly",
                "Builds the [N x T x F] input tensor and graph adjacency matrices consumed by Step 09 training")

    # Dimensions panel
    section_head(s, "Cube dimensions", Inches(0.35), Inches(1.48), w=Inches(3.6))
    for i, (dim, val, desc) in enumerate([
        ("N", "1,397 streets",  "node_index order (Arterial + Council Major)"),
        ("T", "14,400 bins",    "15-min resolution, Nov 2025 - Mar 2026"),
        ("F", "23 features",    "2 dynamic signals + time + weather + 9 static"),
    ]):
        bx = Inches(0.35 + i * 4.3)
        rect(s, bx, Inches(1.95), Inches(3.9), Inches(1.05), NAVY if i == 0 else BLUE if i == 1 else BLUE_LIGHT)
        txbox(s, dim,  bx + Inches(0.14), Inches(1.99), Inches(0.55), Inches(0.5),
              size=22, bold=True, color=WHITE)
        txbox(s, val,  bx + Inches(0.72), Inches(2.0),  Inches(2.9),  Inches(0.38),
              size=13, bold=True, color=WHITE)
        txbox(s, desc, bx + Inches(0.72), Inches(2.44), Inches(2.9),  Inches(0.38),
              size=9.5, color=WHITE)

    rect(s, Inches(0.35), Inches(3.1), Inches(12.98), Inches(0.025), GREY_BOX)

    # Feature table (left column)
    section_head(s, "23 Feature layout", Inches(0.35), Inches(3.22), w=Inches(6.3))
    rows = [
        ("0",   "ped_flow",            "time-varying, normalised"),
        ("1",   "occupancy_rate",      "time-varying, normalised"),
        ("2",   "ped_confidence",      "1.0 / 0.8 / 0.5 — NOT normalised"),
        ("3-6", "hour/dow sin+cos",    "cyclic [-1, 1] — NOT normalised"),
        ("7-9", "is_weekend/holiday",  "binary — NOT normalised"),
        ("10-13","temperature/humidity","time-varying, normalised"),
        ("14-22","total_jobs, cafes,",  "STATIC broadcast across T, normalised"),
        ("",    "bars, business, ...",  ""),
    ]
    for i, (idx, name, note) in enumerate(rows):
        by = Inches(3.68 + i * 0.43)
        fill = GREY_LIGHT if i % 2 == 0 else WHITE
        rect(s, Inches(0.35), by, Inches(6.3), Inches(0.42), fill)
        txbox(s, idx,  Inches(0.42), by + Inches(0.07), Inches(0.6),  Inches(0.3), size=9,  bold=True, color=NAVY)
        txbox(s, name, Inches(1.05), by + Inches(0.07), Inches(2.5),  Inches(0.3), size=9,  color=BLACK)
        txbox(s, note, Inches(3.65), by + Inches(0.07), Inches(3.0),  Inches(0.3), size=8.5, color=GREY_MID)

    # Right panel: adjacency + normalisation
    section_head(s, "Graph adjacency matrices", Inches(6.85), Inches(3.22), w=Inches(6.1))
    code_box(s, [
        "A_norm = D^{-1/2} (A + I) D^{-1/2}",
        "",
        "A       = raw edge weights (Gaussian / cosine ratio)",
        "I       = identity (self-loops: every node gets its own signal)",
        "D       = degree matrix (row sums of A + I)",
        "",
        "Stored as sparse COO torch tensors (.pt)",
        "  -> graph_spatial.pt   (~37,670 bidirectional edges)",
        "  -> graph_semantic.pt  (~9,852  bidirectional edges)",
    ], Inches(6.85), Inches(3.7), Inches(6.1), Inches(2.52))

    section_head(s, "Normalisation statistics", Inches(6.85), Inches(6.38), w=Inches(6.1))
    bullet_lines(s, [
        "Computed on TRAIN portion only (first 80% of timesteps) -- no leakage into val/test",
        "Mean + std for each of the 15 normalisable features  ->  norm_stats.json",
        "Applied on-the-fly during training; NOT baked into cube.npy",
    ], Inches(6.85), Inches(6.82), Inches(6.1), Inches(0.6), size=10.5, color=GREY_DARK)

    footer_note(s, "cube.npy is stored as raw float32 (~0.5-2 GB). Normalisation (z-score) is applied per-batch during Step 09 training using norm_stats.json.")


def s_step9_arch(prs):
    """Step 9 — MultiGCN Architecture."""
    s = blank_slide(prs)
    bg(s, WHITE)
    top_bar(s, "STEP 09  |  TRAIN  |  model architecture", NAVY)
    slide_title(s, "MultiGCN Architecture",
                "GRU encoder + dual GCNConv branches (spatial + semantic) + per-node bias")

    # Architecture diagram (text-based)
    code_box(s, [
        "Input:  [B, W, N, F]    B=batch  W=96 timesteps  N=1,397 streets  F=23 features",
        "",
        "  Spatial branch:   proj_s   (F=23 -> H=64)  ->  ReLU( adj_spatial  @ proj_s(x)  )   [K, N, 64]",
        "  Semantic branch:  proj_sem (F=23 -> H=64)  ->  ReLU( adj_semantic @ proj_sem(x))   [K, N, 64]",
        "",
        "  Concatenate both branches:  [K, N, 128]",
        "  Reshape to:                 [B, W, N, 128]",
        "  Dropout(p=0.10)",
        "",
        "  GRU:  input=128, hidden=64, layers=2, batch_first=True",
        "        Applied per-node:  [B x N, W, 128]  ->  final hidden state  [B, N, 64]",
        "",
        "  Linear:  64  ->  1",
        "  + node_bias [N, 1]   (one learned offset per street)",
        "",
        "Output: [B, N, 1]   (predicted ped_flow, z-score normalised)",
    ], Inches(0.35), Inches(1.42), Inches(8.3), Inches(5.4))

    # Key design choices panel
    section_head(s, "Key design choices", Inches(8.9), Inches(1.42), Inches(4.1))
    bullet_lines(s, [
        "Dual GCN runs at EVERY timestep",
        "  -> GRU gets spatially-aware input",
        "  -> not just spatial features at start",
        "",
        "Per-node bias = highest single",
        "  improvement for R2 on heterogeneous",
        "  streets (busy vs. quiet)",
        "",
        "GRU compresses W=96 steps (24h)",
        "  into a single hidden state",
        "",
        "~100,000 trainable parameters",
    ], Inches(8.9), Inches(1.88), Inches(4.1), Inches(5.0),
       size=10.5, color=GREY_DARK)


def s_step9_training(prs):
    """Step 9 — Training Protocol."""
    s = blank_slide(prs)
    bg(s, WHITE)
    top_bar(s, "STEP 09  |  TRAIN  |  training protocol", NAVY)
    slide_title(s, "Training Protocol",
                "Chronological split, MAE loss, early stopping on validation MAE")

    # Data split
    section_head(s, "Data split (chronological, no leakage)", Inches(0.35), Inches(1.45))
    rect(s, Inches(0.35), Inches(1.95), Inches(8.4), Inches(0.7), RGBColor(0xD0, 0xE4, 0xFF))
    rect(s, Inches(8.75), Inches(1.95), Inches(1.8), Inches(0.7), RGBColor(0xFF, 0xE6, 0xBB))
    rect(s, Inches(10.55), Inches(1.95), Inches(2.43), Inches(0.7), RGBColor(0xD4, 0xED, 0xDA))
    txbox(s, "TRAIN  (70%)   bins 0 - 10,079   Nov 2025 - early Feb 2026",
          Inches(0.5), Inches(2.05), Inches(8.0), Inches(0.42), size=10, color=NAVY)
    txbox(s, "VAL  (15%)", Inches(8.8), Inches(2.05),
          Inches(1.6), Inches(0.42), size=10, color=RGBColor(0x8B, 0x60, 0x00))
    txbox(s, "TEST  (15%)  held-out",
          Inches(10.58), Inches(2.05), Inches(2.3), Inches(0.42), size=10,
          color=RGBColor(0x1A, 0x7A, 0x4A))

    # Hyperparameters grid
    section_head(s, "Hyperparameters", Inches(0.35), Inches(2.85))
    params = [
        ("Batch size",      "8 windows"),
        ("Steps / epoch",   "256 gradient steps"),
        ("Window W",        "96 bins  (24 hours)"),
        ("Loss",            "MAE on z-score target"),
        ("Optimiser",       "Adam  lr = 1e-3"),
        ("LR scheduler",    "ReduceLROnPlateau\nfactor=0.5, patience=8"),
        ("Grad clip",       "max norm 1.0"),
        ("Early stopping",  "patience=25 epochs\non val MAE"),
        ("Seed",            "42 (fully reproducible)"),
        ("Device",          "CUDA if available"),
    ]
    for i, (k, v) in enumerate(params):
        col = i % 5
        row = i // 5
        bx = Inches(0.35 + col * 2.6)
        by = Inches(3.32 + row * 1.2)
        info_box(s, k, v, bx, by, w=Inches(2.48), h=Inches(1.05))

    # Outputs
    section_head(s, "Outputs", Inches(0.35), Inches(5.72), w=Inches(12.6))
    bullet_lines(s, [
        "data/models/best_model.pt  -  trained MultiGCN checkpoint (weights at best val MAE)",
        "data/models/model_eval.json  -  val MAE, val RMSE, val R2, test MAE, test RMSE, test R2, full training loss history",
        "data/processed/cube_meta.json  -  updated with T_train_end, T_val_start, T_val_end indices",
    ], Inches(0.35), Inches(6.18), Inches(12.6), Inches(1.1), size=11, color=GREY_DARK)


def s_step10(prs):
    """Step 10 — Feature Interpretation."""
    s = blank_slide(prs)
    bg(s, WHITE)
    top_bar(s, "STEP 10  |  INTERPRET  |  permutation importance + branch contribution", BLUE)
    slide_title(s, "Feature Interpretation",
                "Which of the 23 features does the model actually rely on? How much does each graph branch contribute?")

    # Left: permutation importance
    section_head(s, "Permutation Importance (23 features)", Inches(0.35), Inches(1.45), w=Inches(6.0))
    bullet_lines(s, [
        "For each feature i = 0 .. 22:",
        "   1. Load val-period windows",
        "   2. Replace feature i with N(0,1) Gaussian noise (in normalised input space)",
        "   3. Measure MAE on perturbed data",
        "   4. importance_i = perturbed_MAE - baseline_MAE",
        "      (positive = model relied on this feature;",
        "       destroying it increases error)",
        "   5. Repeat 3 times, report mean delta-MAE",
        "",
        "Features sorted descending by importance.",
        "Features near 0 delta are not used by the model.",
    ], Inches(0.35), Inches(1.88), Inches(6.1), Inches(4.0), size=11, color=GREY_DARK)

    divider(s, Inches(1.45), GREY_LIGHT)
    rect(s, Inches(6.55), Inches(1.42), Inches(0.03), Inches(5.7), GREY_BOX)

    # Right: branch contribution
    section_head(s, "Branch Contribution", Inches(6.8), Inches(1.45), w=Inches(6.1))
    bullet_lines(s, [
        "For each GCN branch (spatial, semantic):",
        "   1. Zero all weights of that branch's projection layer",
        "      (proj_s or proj_sem)",
        "   2. Measure MAE with the branch disabled",
        "   3. branch_delta_MAE = disabled_MAE - baseline_MAE",
        "   4. Restore original weights",
        "",
        "Interpretation:",
        "   Large delta  ->  branch carries critical information",
        "   Near-zero    ->  branch redundant (other branch compensates)",
        "",
        "Spatial branch  = local neighbourhood context",
        "Semantic branch = functionally similar streets (same land use)",
    ], Inches(6.8), Inches(1.88), Inches(6.0), Inches(4.0), size=11, color=GREY_DARK)

    # Outputs footer
    rect(s, Inches(0.35), Inches(6.1), Inches(12.6), Inches(1.0), GREY_LIGHT)
    txbox(s, "Output: data/models/feature_importance.json",
          Inches(0.55), Inches(6.18), Inches(6.0), Inches(0.35),
          size=10, bold=True, color=NAVY)
    bullet_lines(s, [
        "baseline_mae  |  feature_importance: { feature_name: delta_MAE, ... } sorted descending  |  branch_contribution: { spatial: delta, semantic: delta }",
    ], Inches(0.55), Inches(6.5), Inches(12.2), Inches(0.45), size=10, color=GREY_DARK)


def s_step11_overview(prs):
    """Step 11 — Scenario Simulation overview."""
    s = blank_slide(prs)
    bg(s, WHITE)
    top_bar(s, "STEP 11  |  SCENARIO  |  counterfactual simulation", BLUE_LIGHT)
    slide_title(s, "Scenario Simulation",
                "Predict pedestrian flow impact of a parking intervention on a specific street + time window")

    # 3 intervention types
    section_head(s, "Intervention types", Inches(0.35), Inches(1.45), w=Inches(3.8))
    for i, (name, desc, col) in enumerate([
        ("pedestrianise", "Set target street occupancy_rate = 0\n(full removal of parking)", RGBColor(0xC0, 0x39, 0x2B)),
        ("restrict_park", "Set occupancy_rate = magnitude\n(e.g. 0.3 = allow only 30% use)", BLUE),
        ("boost_ped",     "Add a raw pedestrian uplift\n(ped / 15-min) to predicted flow", ACCENT_GRN),
    ]):
        bx = Inches(0.35 + i * 4.3)
        rect(s, bx, Inches(1.9), Inches(4.0), Inches(1.3), col)
        txbox(s, name, bx + Inches(0.12), Inches(1.96),
              Inches(3.8), Inches(0.42), size=11, bold=True, color=WHITE)
        txbox(s, desc, bx + Inches(0.12), Inches(2.34),
              Inches(3.8), Inches(0.72), size=10, color=WHITE)

    # Rollout logic
    section_head(s, "Autoregressive rollout logic", Inches(0.35), Inches(3.38), w=Inches(6.4))
    code_box(s, [
        "Initialise: sliding window [W=96 real steps before t_start]",
        "",
        "For each rollout step k:",
        "  1. Forward pass  ->  model predicts ped_flow for all 1,397 streets",
        "  2. Build next_row from real observed features (weather, time, occ, static)",
        "  3. Override ped_flow with model prediction  (autoregressive feedback)",
        "  4. If k is inside intervention window:",
        "       apply intervention encoding on target street",
        "  5. If parking removed AND adj_spatial available:",
        "       redistribute displaced occupancy to spatial neighbours",
        "  6. Append next_row, pop oldest from window",
        "",
        "Run TWICE: baseline (no change) and treated (with intervention)",
        "delta = treated - baseline  at every (step, street) pair",
    ], Inches(0.35), Inches(3.85), Inches(6.55), Inches(3.42))

    # Two-rollout diagram
    section_head(s, "Two-rollout design", Inches(7.15), Inches(3.38), w=Inches(5.85))
    bullet_lines(s, [
        "BASELINE rollout: real data, no change",
        "   -> predicted ped_flow for all streets",
        "",
        "TREATED rollout: same data +",
        "   intervention applied to target street",
        "   -> counterfactual ped_flow",
        "",
        "delta = treated - baseline",
        "   -> isolates the intervention effect",
        "   -> controls for weather, time-of-day,",
        "      and seasonal trends",
        "",
        "Both rollouts use same real features.",
        "Only the intervention encoding differs.",
    ], Inches(7.15), Inches(3.85), Inches(5.85), Inches(3.42), size=11, color=GREY_DARK)


def s_step11_network(prs):
    """Step 11 — 5 Network Analysis Enhancements."""
    s = blank_slide(prs)
    bg(s, WHITE)
    top_bar(s, "STEP 11  |  SCENARIO  |  network spillover analysis", BLUE_LIGHT)
    slide_title(s, "5 Network Analysis Enhancements",
                "How intervention effects propagate across 1,397 streets")

    enhancements = [
        (
            "1  Parking Displacement",
            NAVY,
            [
                "When parking is removed from the target street,",
                "displaced occupancy is redistributed to spatial",
                "neighbours proportional to edge weight.",
                "Capped at 100% occupancy per neighbour.",
                "-> Neighbouring streets see realistic input changes.",
            ],
        ),
        (
            "2  Graph Diffusion Analysis",
            BLUE,
            [
                "Analytically propagates mean_delta through",
                "row-normalised adjacency matrices:",
                "  k=1, 2, 3 hops  (A^k @ mean_delta)",
                "Reports top-5 streets per hop for both",
                "spatial and semantic graphs.",
                "-> Complements autoregressive estimate.",
            ],
        ),
        (
            "3  Semantic Neighbour Reporting",
            BLUE_LIGHT,
            [
                "Reports delta on functionally similar streets",
                "(same land-use profile, semantic graph edges).",
                "These streets may be spatially distant.",
                "-> Captures non-local demand shifts where",
                "   displaced users travel to similar streets.",
            ],
        ),
        (
            "4  Confidence-Weighted Ranking",
            RGBColor(0x70, 0x50, 0xC0),
            [
                "top_affected streets ranked by:",
                "  delta x ped_confidence",
                "Sensor streets (conf=1.0) ranked above",
                "imputed streets (conf=0.5) for same delta.",
                "-> Avoids promoting noisy imputed streets",
                "   over reliable sensor observations.",
            ],
        ),
        (
            "5  Rebound / Half-Life Analysis",
            RGBColor(0x1A, 0x7A, 0x4A),
            [
                "After intervention ends, measures:",
                "  - Steps until |delta| drops to 50%",
                "    of its end-of-intervention value",
                "  - Recovery fraction at final rollout step",
                "Per-street half-life in steps and minutes.",
                "-> Characterises persistence of effect.",
            ],
        ),
    ]

    for i, (title, col, bullets) in enumerate(enhancements):
        col_idx = i % 3
        row_idx = i // 3
        bx = Inches(0.35 + col_idx * 4.34)
        by = Inches(1.48 + row_idx * 2.85)
        bw = Inches(4.1)
        bh = Inches(2.62)
        rect(s, bx, by, bw, bh, GREY_LIGHT)
        rect(s, bx, by, bw, Inches(0.4), col)
        txbox(s, title, bx + Inches(0.1), by + Inches(0.05),
              bw - Inches(0.15), Inches(0.34),
              size=10.5, bold=True, color=WHITE)
        bullet_lines(s, bullets, bx + Inches(0.15), by + Inches(0.48),
                     bw - Inches(0.25), bh - Inches(0.58),
                     size=10, color=GREY_DARK)


def s_step11_output(prs):
    """Step 11 — Output JSON structure."""
    s = blank_slide(prs)
    bg(s, WHITE)
    top_bar(s, "STEP 11  |  SCENARIO  |  output structure + usage", BLUE_LIGHT)
    slide_title(s, "Scenario Output Structure",
                "Saved as JSON; also returned live via the Flask API (Step 11 API server)")

    code_box(s, [
        '{',
        '  "meta":    { street_id, t_start_bin, duration_bins, rollout_steps, intervention_type }',
        '  "baseline":{ ped_flow_treated_street, occ_rate_observed }',
        '  "treated": { ped_flow_treated_street }',
        '  "delta":   { ped_flow_treated_street, top_affected_series }',
        '  "network_summary": {',
        '     "treated_street":      { mean_delta, cumulative_delta, ped_confidence }',
        '     "spatial_neighbours":  [ { street_id, mean_delta, confidence_weighted_delta } ]',
        '     "semantic_neighbours": [ { street_id, mean_delta, confidence_weighted_delta } ]',
        '     "top_affected":        [ { street_id, mean_delta, is_treated,',
        '                               is_spatial_neighbour, is_semantic_neighbour } ]',
        '     "graph_diffusion":     { "spatial": [...k=1,2,3 top-5...],',
        '                             "semantic": [...k=1,2,3 top-5...] }',
        '     "rebound":             { "treated_street": { half_life_steps, peak_delta,',
        '                               recovery_fraction },',
        '                             "slowest_to_recover": [...] }',
        '     "all_deltas":          { street_id: { mean_ped_delta,',
        '                               confidence_weighted_delta, ped_confidence } }',
        '  }',
        '}',
    ], Inches(0.35), Inches(1.5), Inches(7.65), Inches(5.72))

    # CLI + API usage
    section_head(s, "CLI usage", Inches(8.2), Inches(1.5), w=Inches(4.8))
    code_box(s, [
        "python melbourne_pipeline/steps/",
        "  step_11_scenario.py \\",
        "  --street STREET_ID \\",
        "  --start BIN_INDEX \\",
        "  --duration BINS \\",
        "  --rollout STEPS \\",
        "  --intervention pedestrianise",
    ], Inches(8.2), Inches(1.95), Inches(4.8), Inches(2.15))

    section_head(s, "API (Flask, port 5050)", Inches(8.2), Inches(4.3), w=Inches(4.8))
    code_box(s, [
        "POST http://127.0.0.1:5050/scenario",
        "",
        "{",
        '  "street_id": "21373",',
        '  "intervention_type": "pedestrianise",',
        '  "duration": 16,',
        '  "rollout_steps": 32',
        "}",
    ], Inches(8.2), Inches(4.75), Inches(4.8), Inches(2.0))

    footer_note(s, "Recommended test streets: 21373 (highest pedestrian flow), 21123, 22619  |  Use high-occupancy streets to observe meaningful deltas")


def s_step11b(prs):
    """Step 11b — Opportunity Scoring."""
    s = blank_slide(prs)
    bg(s, WHITE)
    top_bar(s, "STEP 11b  |  OPPORTUNITIES  |  composite intervention ranking", BLUE_LIGHT)
    slide_title(s, "Opportunity Scoring",
                "Ranks all 1,397 streets by their suitability for curbside reallocation intervention")

    section_head(s, "Composite Opportunity Score  (0 - 1)", Inches(0.35), Inches(1.48), w=Inches(7.0))
    code_box(s, [
        "score = 0.35 x score_ped        (predicted pedestrian flow, normalised by max)",
        "      + 0.25 x score_parking    (1 - mean_occupancy: low parking = more space)",
        "      + 0.20 x score_archetype  (GMM cluster confidence from Step 07)",
        "      + 0.10 x score_confidence (data source: sensor=1.0 / imputed R2>=0.6=0.8)",
        "      + 0.10 x score_uplift     (counterfactual uplift fraction)",
    ], Inches(0.35), Inches(1.95), Inches(12.6), Inches(1.72))

    # Weight breakdown visual
    weights = [
        ("Pedestrian\ndemand", 0.35, NAVY),
        ("Parking\navailability", 0.25, BLUE),
        ("Archetype\nconfidence", 0.20, BLUE_LIGHT),
        ("Data\nconfidence", 0.10, GREY_MID),
        ("Predicted\nuplift", 0.10, GREY_MID),
    ]
    total_w = Inches(12.6)
    x_cur = Inches(0.35)
    bar_y = Inches(3.85)
    bar_h = Inches(0.6)
    for label, wt, col in weights:
        bw = total_w * wt
        rect(s, x_cur, bar_y, bw, bar_h, col)
        pct = f"{int(wt*100)}%"
        txbox(s, pct, x_cur + Inches(0.05), bar_y + Inches(0.1),
              bw - Inches(0.08), Inches(0.4), size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        txbox(s, label, x_cur, bar_y + Inches(0.68),
              bw, Inches(0.5), size=8.5, color=GREY_DARK, align=PP_ALIGN.CENTER)
        x_cur += bw

    section_head(s, "Counterfactual uplift computation", Inches(0.35), Inches(4.72), w=Inches(6.5))
    bullet_lines(s, [
        "Run 8 randomly sampled val-period windows",
        "Set occupancy_rate = 0 for ALL streets simultaneously",
        "Measure mean predicted pedestrian uplift vs. baseline",
        "uplift_fraction = (cf_ped_mean - pred_ped_mean) / (pred_ped_mean + 1e-6)",
    ], Inches(0.35), Inches(5.16), Inches(6.5), Inches(1.5), size=11, color=GREY_DARK)

    section_head(s, "Outputs", Inches(7.1), Inches(4.72), w=Inches(5.9))
    bullet_lines(s, [
        "street_scores.parquet",
        "   street_id, opportunity_score, pred_ped_mean,",
        "   uplift_fraction, cf_ped_mean, cluster, ...",
        "",
        "opportunities.json",
        "   Top-ranked streets with all score components",
    ], Inches(7.1), Inches(5.16), Inches(5.9), Inches(1.5), size=11, color=GREY_DARK)


def s_step12(prs):
    """Step 12 — Frontend Export."""
    s = blank_slide(prs)
    bg(s, WHITE)
    top_bar(s, "STEP 12  |  EXPORT  |  GeoJSON + JSON for Mapbox GL JS", GREY_DARK)
    slide_title(s, "Frontend Export",
                "Enriches spatial data with pipeline outputs ready for sensor_map_viz.html")

    section_head(s, "Operations", Inches(0.35), Inches(1.48), w=Inches(12.6))
    ops = [
        ("1", "Join opportunity scores",
         "Loads street_scores.parquet. Joins: opportunity_score, pred_ped_mean, uplift_fraction,\ncf_ped_mean onto each feature in streets_viz.geojson (in-place enrichment)."),
        ("2", "Write opportunities summary",
         "Writes opportunities_summary.json -- top-50 streets per archetype with key metrics.\nUsed by the frontend dropdown filters."),
        ("3", "Copy model evaluation",
         "Copies model_eval.json from data/models/ to data/processed/ so the frontend\ncan read val/test MAE, RMSE, R2 without knowing the models directory."),
    ]
    for i, (num, title, desc) in enumerate(ops):
        by = Inches(1.95 + i * 1.45)
        rect(s, Inches(0.35), by, Inches(0.55), Inches(1.25), NAVY)
        txbox(s, num, Inches(0.35), by + Inches(0.38), Inches(0.55), Inches(0.5),
              size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        rect(s, Inches(0.9), by, Inches(12.08), Inches(1.25), GREY_LIGHT)
        txbox(s, title, Inches(1.06), by + Inches(0.1), Inches(11.7), Inches(0.38),
              size=12, bold=True, color=NAVY)
        txbox(s, desc, Inches(1.06), by + Inches(0.48), Inches(11.7), Inches(0.65),
              size=10.5, color=GREY_DARK)

    # Output files
    section_head(s, "Output files", Inches(0.35), Inches(6.27), w=Inches(12.6))
    bullet_lines(s, [
        "data/processed/streets_viz.geojson         Enriched street polygons (WGS84) with all pipeline scores for Mapbox rendering",
        "data/processed/opportunities_summary.json  Top-50 streets per archetype -- cluster label, opportunity_score, pred_ped_mean, uplift_fraction",
        "data/processed/model_eval.json             Copied model performance metrics  (val/test MAE, RMSE, R2, training history)",
    ], Inches(0.35), Inches(6.7), Inches(12.6), Inches(0.65), size=10.5, color=GREY_DARK)


def s_summary(prs):
    """Summary slide: data flow + key facts."""
    s = blank_slide(prs)
    bg(s, WHITE)
    top_bar(s, "SUMMARY  |  Steps 9 - 12 data flow and key facts", NAVY)
    slide_title(s, "Summary: Steps 9 - 12",
                "From trained weights to interactive intervention explorer")

    # Data flow diagram
    section_head(s, "Data flow", Inches(0.35), Inches(1.48), w=Inches(6.5))

    flow = [
        ("cube.npy\ngraph_spatial.pt\ngraph_semantic.pt",  GREY_BOX,  GREY_DARK),
        ("Step 09\nTRAIN",                                  NAVY,      WHITE),
        ("best_model.pt\nmodel_eval.json",                  GREY_BOX,  GREY_DARK),
        ("Step 10\nINTERPRET",                              BLUE,      WHITE),
        ("feature_importance.json",                         GREY_BOX,  GREY_DARK),
    ]
    flow2 = [
        ("best_model.pt\ncube.npy",                         GREY_BOX,  GREY_DARK),
        ("Step 11\nSCENARIO",                               BLUE_LIGHT, WHITE),
        ("scenario JSON\nall_deltas",                       GREY_BOX,  GREY_DARK),
        ("Step 12\nEXPORT",                                 GREY_DARK, WHITE),
        ("streets_viz.geojson\nopportunities_summary.json", GREY_BOX,  GREY_DARK),
    ]
    for row_i, flow_row in enumerate([flow, flow2]):
        for i, (label, fill, txt) in enumerate(flow_row):
            bx = Inches(0.35 + i * 2.55)
            by = Inches(1.98 + row_i * 2.0)
            bw = Inches(2.3)
            bh = Inches(0.85)
            rect(s, bx, by, bw, bh, fill)
            txbox(s, label, bx + Inches(0.08), by + Inches(0.1),
                  bw - Inches(0.12), bh - Inches(0.12),
                  size=9, bold=(fill not in (GREY_BOX,)), color=txt, align=PP_ALIGN.CENTER)
            if i < 4:
                txbox(s, "->", bx + bw, by + Inches(0.28), Inches(0.25), Inches(0.35),
                      size=12, bold=True, color=GREY_MID, align=PP_ALIGN.CENTER)

    # Key facts grid
    section_head(s, "Key numbers", Inches(0.35), Inches(5.98), w=Inches(12.6))
    facts = [
        ("Model params",    "~100,000"),
        ("Training windows","~2,560 / epoch"),
        ("Inference time",  "~200 ms per scenario"),
        ("Network coverage","1,397 streets"),
        ("Rollout steps",   "32 - 96 typical"),
        ("Confidence tiers","1.0 / 0.8 / 0.5"),
    ]
    for i, (k, v) in enumerate(facts):
        info_box(s, k, v,
                 Inches(0.35 + i * 2.17), Inches(6.45),
                 w=Inches(2.0), h=Inches(0.8))


# ═══════════════════════════════════════════════════════════════════════════════
# BUILD
# ═══════════════════════════════════════════════════════════════════════════════

def build():
    prs = new_prs()
    s_cover(prs)
    s_overview(prs)
    s_step8(prs)
    s_step9_arch(prs)
    s_step9_training(prs)
    s_step10(prs)
    s_step11_overview(prs)
    s_step11_network(prs)
    s_step11_output(prs)
    s_step11b(prs)
    s_step12(prs)
    s_summary(prs)
    prs.save(str(OUT_PATH))
    print(f"Saved -> {OUT_PATH}  ({len(prs.slides)} slides)")


if __name__ == "__main__":
    build()
